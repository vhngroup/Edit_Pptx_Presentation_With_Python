from pptx import Presentation

file_path = 'base.pptx'
wilcard_Replacement = {
    "%TITLE%": "VHNGROUP",
    "%subtitulo%":"Integramos Seguridad y Tecnologia \n Plan Gerencia",
    "%TITULO1%":"Presentación de Producto",
    "%texto_Interno%":"Demo de texto para Presentation, gerencial"
}
rectangle_replament = {
    "IMG1":"image1.jpg",
    "IMG2":"image2.jpg"
}

def process_Pptx(file, wilcard,rectangle ):
    ppt = Presentation(file)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.shape_type ==1:
                replace_rect_img(slide, shape, rectangle)
            elif hasattr(shape, 'text'):
                repace_text(shape, wilcard) 
    ppt.save('Salida.pptx')
    print("Archivo creado satidfactoriamente")

def repace_text(shape, wilcard):
    if shape.text in wilcard:
        shape.text = wilcard[shape.text]

def replace_rect_img(slide, shape, rectangle):
    if shape.text in rectangle:
        left = shape.left
        top = shape.top
        width= shape.width
        height= shape.height

        slide.shapes._spTree.remove(shape._element)
        slide.shapes.add_picture(rectangle[shape.text], left, top, width, height)

process_Pptx(file_path, wilcard_Replacement, rectangle_replament)